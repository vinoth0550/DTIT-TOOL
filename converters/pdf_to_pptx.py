
# converters/pdf_to_ppt.py

from fastapi import APIRouter, File, UploadFile, HTTPException
from fastapi.responses import FileResponse, JSONResponse
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import fitz, tempfile, uuid, io

from utils.storage_manager import save_upload, build_output_path
import config

router = APIRouter(
    prefix="/pdf-to-ppt",
    tags=["PDF → PPT"]
)

TOOL_NAME = "pdf_to_ppt"


# ORIGINAL LOGIC 

def pdf_to_images(pdf_bytes: bytes) -> list:
    images = []
    pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")

    for page_num in range(pdf_document.page_count):
        page = pdf_document[page_num]
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        images.append(pix.tobytes("png"))

    pdf_document.close()
    return images


def create_ppt(images: list, output_path: str):
    prs = Presentation()
    
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    for img_bytes in images:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        img = Image.open(io.BytesIO(img_bytes))
        width, height = img.size

        slide_aspect = prs.slide_width / prs.slide_height
        img_aspect = width / height

        if img_aspect > slide_aspect:
            new_width = prs.slide_width
            new_height = prs.slide_width / img_aspect
        else:
            new_height = prs.slide_height
            new_width = prs.slide_height * img_aspect

        left = (prs.slide_width - new_width) / 2
        top = (prs.slide_height - new_height) / 2

        temp_img = tempfile.gettempdir() + f"/{uuid.uuid4()}.png"
        with open(temp_img, "wb") as f:
            f.write(img_bytes)

        slide.shapes.add_picture(temp_img, left, top, width=new_width, height=new_height)

    prs.save(output_path)


# API ROUTE

@router.post("/")
async def convert_pdf(file: UploadFile = File(...)):

    if not file.filename.lower().endswith(".pdf"):
        return JSONResponse(status_code=400, content={"status": "error", "message": "Only PDF files allowed"})

    pdf_bytes = await file.read()
    if len(pdf_bytes) == 0:
        return JSONResponse(status_code=400, content={"status": "error", "message": "Empty PDF file"})

    try:
        # CREATE OUTPUT WITH COUNTER LOGIC
        output_path = build_output_path("converted-slide", ".pptx", TOOL_NAME)

        # ORIGINAL CONVERSION LOGIC 
        images = pdf_to_images(pdf_bytes)
        if not images:
            return JSONResponse(status_code=400, content={"status": "error", "message": "PDF has no pages"})

        create_ppt(images, str(output_path))

        download_url = f"{config.BASE_DOWNLOAD_URL}/{TOOL_NAME}/{output_path.name}"

        return {
            "status": "success",
            "message": "Converted PDF → PPT successfully!",
            "download_link": download_url,
            "file_name": output_path.name
        }

    except Exception as e:
        return JSONResponse(status_code=500, content={"status": "error", "message": str(e)})


@router.get("/file/{filename}")
def download(filename: str):
    file_path = config.OUTPUT_ROOT / TOOL_NAME / filename

    if not file_path.exists():
        return JSONResponse(status_code=404, content={"status": "error", "message": "File not found"})

    return FileResponse(
        path=str(file_path),
        filename=filename,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
