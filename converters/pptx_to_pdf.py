

# converters/pptx_to_pdf.py

from fastapi import APIRouter, UploadFile, File, HTTPException
from fastapi.responses import FileResponse, JSONResponse
import uuid, os, platform, subprocess, tempfile
from pathlib import Path

from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import landscape
from PIL import Image
import fitz

from utils.storage_manager import save_upload, build_output_path
import config








router = APIRouter(
    prefix="/pptx-to-pdf",
    tags=["PPTX → PDF"]
)

TOOL_NAME = "pptx_to_pdf"


# ---------------- ORIGINAL LOGIC (NOT MODIFIED) ---------------- #

def pptx_to_pdf_windows(pptx_path: str, pdf_path: str):
    import comtypes.client
    try:
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1
        deck = powerpoint.Presentations.Open(pptx_path)
        deck.SaveAs(pdf_path, 32)
        deck.Close()
        powerpoint.Quit()
        return True
    except:
        return False


def pptx_to_pdf_libreoffice(pptx_path: str, pdf_path: str):
    try:
        commands = ["libreoffice", "soffice", "/usr/bin/libreoffice"]
        for cmd in commands:
            try:
                result = subprocess.run([
                    cmd, "--headless", "--convert-to", "pdf",
                    "--outdir", str(Path(pdf_path).parent), pptx_path
                ], capture_output=True, timeout=50)

                if result.returncode == 0:
                    generated = Path(pptx_path).with_suffix(".pdf")
                    if generated.exists():
                        if generated != pdf_path:
                            os.replace(generated, pdf_path)
                        return True
            except:
                continue
        return False
    except:
        return False


def pptx_to_pdf_fallback(pptx_path: str, pdf_path: str):
    try:
        prs = Presentation(pptx_path)
        slide_width = prs.slide_width.inches
        slide_height = prs.slide_height.inches

        c = canvas.Canvas(pdf_path, pagesize=landscape((slide_width * 72, slide_height * 72)))


        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    c.drawString(50, slide_height * 72 - 50, shape.text[:120])
            c.showPage()

        c.save()
        return True
    except:
        return False


def convert_ppt_logic(pptx_path: str, pdf_path: str) -> bool:
    if platform.system() == "Windows" and pptx_to_pdf_windows(pptx_path, pdf_path):
        return True

    if pptx_to_pdf_libreoffice(pptx_path, pdf_path):
        return True

    return pptx_to_pdf_fallback(pptx_path, pdf_path)


#  API ENDPOINT 

@router.post("/")
async def convert(file: UploadFile = File(...)):

    if not file.filename.lower().endswith((".pptx", ".ppt")):
        raise HTTPException(status_code=400, detail="Only PPTX/PPT files allowed")

    # Save upload to /uploads/
    input_path = save_upload(file, TOOL_NAME)

    # Build output filename with counter logic
    output_path = build_output_path(Path(file.filename).stem, ".pdf", TOOL_NAME)

    # Perform conversion
    success = convert_ppt_logic(str(input_path), str(output_path))

    if not success or not output_path.exists():
        raise HTTPException(status_code=500, detail="Conversion failed. LibreOffice or PowerPoint required.")

    download_url = f"{config.BASE_DOWNLOAD_URL}/{TOOL_NAME}/{output_path.name}"

    return JSONResponse({
        "status": "success",
        "message": "PPTX converted to PDF successfully!",
        "download_link": download_url,
        "file_name": output_path.name
    })


@router.get("/file/{filename}")
async def download(filename: str):
    file_path = config.OUTPUT_ROOT / TOOL_NAME / filename

    if not file_path.exists():
        raise HTTPException(status_code=404, detail="File not found")

    return FileResponse(
        str(file_path),
        filename=filename,
        media_type="application/pdf"
    )
